# -*- coding: utf-8 -*-
"""WCS 함수단 흐름도 (pptx) - 입고/출고 풀사이클을 실제 함수 호출 단위로 그린다.

근거 소스 (2026-09-05 기준)
  TASK/WCS_TASK_HOST/CSrvWork.cs      상위 수신(9911) : ReadRequest ~ ParseOorR
  TASK/WCS_TASK_HOST/CCliWork.cs      상위 송신(9910) : GetSendData ~ GetJobCompleteReport
  TASK/IO_TASK/CLS/Thread/cThread_SCH.cs         스케줄러 : Thread_Doing 1사이클
  TASK/WCS_TASK_CV_BIN/2_Thread/EQP_THREAD/CvThread.cs   C/V PLC
  TASK/WCS_TASK_CV_BIN/2_Thread/EQP_THREAD/VehThread.cs  S/C · RGV PLC

상태 전이는 cThread_SCH.cs DriveSC 주석(895~896행)을 따른다.
  입고 : 99 -> 10 -> 15 -> 35 -> 39 -> 15 -> 25 -> 29 -> 09
  출고 : 99 -> 20 -> 25 -> 29 -> 15 -> 35 -> 39 -> 15 -> 19 -> 09
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

OUT = r'D:\project\LGLS\Renewal\docs\산출물_20260904'
DATE = '2026-09-05'
NAVY = RGBColor(0x1E, 0x27, 0x61)
BLUE = RGBColor(0x26, 0x6E, 0xBE)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
LIGHT = RGBColor(0xEE, 0xF1, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x59, 0x59, 0x59)
LINE = RGBColor(0xB6, 0xBF, 0xD0)
GREEN = RGBColor(0x2E, 0x86, 0x4B)
AMBER = RGBColor(0xC8, 0x7A, 0x0E)
RED = RGBColor(0xB0, 0x2E, 0x2E)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
FONT = '맑은 고딕'

# 프로세스별 색 : 범례와 각 단계 상자가 같은 색을 쓴다
PROC = {
    'HOST': (BLUE,  RGBColor(0xEC, 0xF2, 0xFB), 'WCS_TASK_HOST'),
    'SCH':  (GREEN, RGBColor(0xEC, 0xF6, 0xEF), 'IO_TASK'),
    'CV':   (AMBER, RGBColor(0xFB, 0xF4, 0xE8), 'WCS_TASK_CV / CvThread'),
    'VEH':  (TEAL,  RGBColor(0xEA, 0xF4, 0xF6), 'WCS_TASK_CV / VehThread'),
    'EQP':  (GRAY,  LIGHT,                      '설비 (PLC)'),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
page = [0]


def tb(slide, x, y, w, h, text, size=12, bold=False, color=GRAY,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = s.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, tuple):
            txt, b, sz, col = ln
        else:
            txt, b, sz, col = ln, bold, size, color
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = col
        r.font.name = FONT
    return s


def box(slide, x, y, w, h, text, fill=WHITE, edge=NAVY, txt=NAVY, size=11,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, bold=True, align=PP_ALIGN.CENTER):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge
    sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.07)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text if isinstance(text, list) else [text]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, tuple):
            t, sz, b, c = ln
        else:
            t, sz, b, c = ln, size, bold, txt
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = FONT
    return sh


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=1.5, dash=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    el = c.line._get_or_add_ln()
    el.append(parse_xml('<a:tailEnd %s type="triangle" w="med" len="med"/>' % nsdecls('a')))
    if dash:
        el.insert(0, parse_xml('<a:prstDash %s val="dash"/>' % nsdecls('a')))
    return c


def footer(slide, title):
    page[0] += 1
    tb(slide, 0.45, 7.05, 7, 0.3, 'LG 화학 1동 WCS  |  ' + title, 9, False, GRAY)
    tb(slide, 12.3, 7.05, 0.6, 0.3, str(page[0]), 9, False, GRAY, PP_ALIGN.RIGHT)


def head(slide, text, sub=None):
    tb(slide, 0.45, 0.28, 12.4, 0.5, text, 24, True, NAVY, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        tb(slide, 0.45, 0.80, 12.4, 0.3, sub, 10.5, False, GRAY)


def legend(slide, y=1.05):
    """프로세스 색 범례"""
    x = 0.45
    for key in ('HOST', 'SCH', 'CV', 'VEH'):
        edge, fill, name = PROC[key]
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(0.20), Inches(0.16))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = edge
        sh.line.width = Pt(1.0)
        sh.shadow.inherit = False
        tb(slide, x + 0.26, y - 0.03, 2.6, 0.24, name, 9, False, GRAY)
        x += 2.95


def step(slide, x, y, w, h, kind, no, fn, desc, badge=None):
    """함수 단계 상자 : 번호 + 함수명 + 설명 (+ 상태 전이 배지)"""
    edge, fill, _ = PROC[kind]
    lines = [('%s  %s' % (no, fn), 9.5, True, NAVY),
             (desc, 8, False, GRAY)]
    if badge:
        lines.append((badge, 8, True, edge))
    return box(slide, x, y, w, h, lines, fill, edge, NAVY, 9.5)


def chain(slide, steps, cols=3, x0=0.45, y0=1.42, w=3.94, h=0.74, gap=0.80, colgap=4.24):
    """단계 목록을 열 단위로 세로 배치하고 열 안에서 아래로 화살표를 잇는다."""
    rows = (len(steps) + cols - 1) // cols
    for i, st in enumerate(steps):
        c, r = divmod(i, rows)
        x = x0 + c * colgap
        y = y0 + r * gap
        step(slide, x, y, w, h, st[0], st[1], st[2], st[3], st[4] if len(st) > 4 else None)
        if r < rows - 1 and i + 1 < len(steps):
            arrow(slide, x + w / 2, y + h, x + w / 2, y + gap, LINE, 1.25)
    return rows


# ══════════════════════════════════════════════════════════════════
# 1. 표지
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
page[0] += 1
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
bg.shadow.inherit = False
tb(s, 1.0, 2.45, 11, 0.6, 'LG 화학 1동 자동창고', 20, False, ICE)
tb(s, 1.0, 3.10, 11, 1.1, 'WCS 함수단 흐름도', 40, True, WHITE)
tb(s, 1.0, 4.30, 11, 0.5, '입고 · 출고 풀사이클을 실제 함수 호출 단위로 정리', 15, False, ICE)
tb(s, 1.0, 4.95, 11, 0.5, 'WCS Renewal (구 ECS 대체)   ·   %s' % DATE, 13, False, ICE)

# ══════════════════════════════════════════════════════════════════
# 2. 읽는 법 · 프로세스별 진입 함수
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
head(s, '읽는 법 · 프로세스별 진입 함수',
     '네 개의 프로세스가 DB(LGLS_MCS_IO)를 매개로 각자 자기 루프를 돈다. 직접 호출하는 관계는 없다.')

box(s, 0.45, 1.35, 2.85, 1.30,
    [('WCS_TASK_HOST', 12, True, NAVY), ('상위 통신', 8.5, False, GRAY),
     ('CSrvWork.ReadRequest()', 9, False, BLUE),
     ('CCliWork.GetSendData()', 9, False, BLUE)],
    PROC['HOST'][1], BLUE, NAVY)
box(s, 3.65, 1.35, 2.85, 1.30,
    [('IO_TASK', 12, True, NAVY), ('스케줄러', 8.5, False, GRAY),
     ('cThread_SCH.Thread_Doing()', 9, False, GREEN),
     ('200ms 주기 · 단일 루프', 8.5, False, GRAY)],
    PROC['SCH'][1], GREEN, NAVY)
box(s, 6.85, 1.35, 2.85, 1.30,
    [('WCS_TASK_CV', 12, True, NAVY), ('C/V 통신', 8.5, False, GRAY),
     ('CvThread.Thread_Doing()', 9, False, AMBER),
     ('15설비 슬롯 순회 · 약 16초', 8.5, False, GRAY)],
    PROC['CV'][1], AMBER, NAVY)
box(s, 10.05, 1.35, 2.85, 1.30,
    [('WCS_TASK_CV', 12, True, NAVY), ('S/C · RGV 통신', 8.5, False, GRAY),
     ('VehThread.Thread_Doing()', 9, False, TEAL),
     ('관측 폴링 + 명령 소비', 8.5, False, GRAY)],
    PROC['VEH'][1], TEAL, NAVY)

for x in (3.30, 6.50, 9.70):
    arrow(s, x, 2.00, x + 0.35, 2.00, LINE, 1.25)
    arrow(s, x + 0.35, 2.20, x, 2.20, LINE, 1.25)

box(s, 3.65, 3.05, 6.05, 0.62,
    [('DB : LGLS_MCS_IO', 12, True, WHITE),
     ('JOB_MST · CV_DATA · SC_DATA_LGLS · RTV_DATA_LGLS', 9, False, WHITE)],
    NAVY, NAVY, WHITE, 12, MSO_SHAPE.FLOWCHART_MAGNETIC_DISK)
for x in (1.87, 5.07, 8.27, 11.47):
    arrow(s, x, 2.65, x, 3.05, LINE, 1.25)

tb(s, 0.45, 3.95, 6.05, 2.9, [
    ('■ 표기 규칙', True, 12, NAVY),
    ('', False, 6, GRAY),
    ('· 상자 안 윗줄 = 실제 함수명. 소스 그대로다.', False, 10.5, GRAY),
    ('· 아랫줄 = 그 함수가 하는 일과 건드리는 테이블/신호.', False, 10.5, GRAY),
    ('· 오른쪽 배지 = JOB_MST.JOB_STATUS 전이. 배지가 없는 단계는', False, 10.5, GRAY),
    ('  상태를 바꾸지 않고 설비 신호만 주고받는다.', False, 10.5, GRAY),
    ('· 상자 색 = 그 함수가 사는 프로세스.', False, 10.5, GRAY),
    ('· 번호는 시간 순서다. 열이 넘어가면 다음 열 맨 위로 이어진다.', False, 10.5, GRAY),
], anchor=MSO_ANCHOR.TOP)

tb(s, 6.85, 3.95, 6.05, 2.9, [
    ('■ 작업 상태 전이 (근거 : cThread_SCH.cs DriveSC 주석)', True, 12, NAVY),
    ('', False, 6, GRAY),
    ('입고   99 → 10 → 15 → 35 → 39 → 15 → 25 → 29 → 09', True, 11, GREEN),
    ('출고   99 → 20 → 25 → 29 → 15 → 35 → 39 → 15 → 19 → 09', True, 11, BLUE),
    ('', False, 6, GRAY),
    ('· 35/39 구간은 RGV(RTV)를 경유할 때만 나타난다.', False, 10, GRAY),
    ('· 19(출고) / 29(입고) 는 상위 완료보고(F)가 소비하는 최종 상태다.', False, 10, GRAY),
    ('  중간 처리가 19/29 를 거치면 조기 완료보고가 나가므로 금지.', False, 10, RED),
    ('· 09 는 완료보고를 끝낸 상태. JOB_MST 에 그대로 남는다', False, 10, GRAY),
    ('  (반자동만 DeleteSemiFinished() 가 삭제한다).', False, 10, GRAY),
], anchor=MSO_ANCHOR.TOP)
footer(s, '읽는 법')

# ══════════════════════════════════════════════════════════════════
# 3. 입고 풀사이클
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
head(s, '입고 풀사이클 — 함수 호출 순서',
     '상위 O 전문 수신부터 F 완료보고까지. JOB_TYP=\'1\', 입고대 124(C/V#12) → S/C 랙 셀 기준.')
legend(s)
chain(s, [
    ('HOST', '01', 'CSrvWork.ReadRequest()',
     'TCP 9911 수신 → CheckHeader() / CheckBody() 로 전문 길이 확인'),
    ('HOST', '02', 'CSrvWork.Parsing() → ParseOorR()',
     'O 전문 고정자리수 파싱 (LuggNo·SrcStn·DestStn·ProductID)'),
    ('HOST', '03', 'IsValidStartStation() / IsValidDestStation()',
     '작업대 번호와 랙 셀 유효성 검증. 실패 시 MakeResponse() 로 반려'),
    ('HOST', '04', 'frmMain.InsertJobMst()',
     'JOB_MST INSERT — INS_USER_ID=\'HOST_TASK\'', 'JOB_STATUS = 99'),
    ('HOST', '05', 'MakeResponse() → SendSock()',
     '상위에 응답 전문(11byte) 회신'),
    ('SCH', '06', 'cThread_SCH.Thread_Doing()',
     '스케줄러 200ms 루프가 신규 작업을 집는다'),
    ('SCH', '07', 'AcceptNewJob()',
     '입고는 C/V 가 먼저 받는다 (START_POS = 1xx 작업대)', '99 → 10'),
    ('SCH', '08', 'DriveCV() → UpdateCvData()',
     'CV_DATA 에 이송 명령 기록 (OD_RQ_YN=\'Y\', DEST_POS_OD)', '10 → 15'),
    ('CV', '09', 'ScanPendingWork() → CvChg_OD_RQ_YN()',
     '쓰기 대기 설비만 골라 XGT PLC 로 이송 명령 WRITE'),
    ('CV', '10', 'CvStatusScenario() → UpdateCvData()',
     'PLC 상태 READ → CV_DATA 미러 (센서·화물번호·에러코드)'),
    ('CV', '11', 'CvEventCheck()',
     'Load / Unload Complete 이벤트 ACK'),
    ('CV', '12', 'CvTrackingWrite()',
     'R 영역에 작업번호 트래킹 WRITE (EncodeJobNoR)'),
    ('SCH', '13', 'DriveRGV() → UpdateRtvData()',
     'RTV 유휴 확인 후 이송 명령. 도착 트랙을 HS_TRACK_NO 에 기록', '15 → 35'),
    ('SCH', '14', 'CompleteRGVReal()',
     'RTV_DATA_LGLS.COMPLETE_RD 소비', '35 → 39'),
    ('SCH', '15', 'LandRgvDrop() → RequestArrivalTrackingWrite()',
     'HS_TRACK_NO 에 화물이 실제로 올라온 것을 보고 다음 구간에 인계', '39 → 15'),
    ('SCH', '16', 'DriveSC() → UpdateScData()',
     '입고는 15 에서 S/C 가 가져간다. SC_DATA_LGLS OD_RQ_YN=\'Y\'', '15 → 25'),
    ('VEH', '17', 'VehThread.ConsumeCommands()',
     'WriteString() / WriteBit() 으로 S/C 에 이송 명령 WRITE'),
    ('VEH', '18', 'VehThread.PollObservations()',
     'S/C 위치·포크·에러 관측 → SC_DATA_LGLS 반영'),
    ('SCH', '19', 'CompleteSC()',
     '크레인이 랙 셀에 넣으면 입고 최종 상태', '25 → 29'),
    ('HOST', '20', 'GetSendData() → GetJobCompleteReport(29)',
     'IsJobExist(29) 로 완료 대상 조회 → F 전문 조립'),
    ('HOST', '21', 'SendSock() → UpdateJobStatusTo(\'09\')',
     '상위에 완료보고(F). 실패하면 29 로 롤백해 다음 주기 재보고', '29 → 09'),
])
footer(s, '입고 풀사이클')

# ══════════════════════════════════════════════════════════════════
# 4. 출고 풀사이클
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
head(s, '출고 풀사이클 — 함수 호출 순서',
     'JOB_TYP=\'2\'. 출고는 S/C 가 먼저 움직인다. 랙 셀 → 출고대 126(C/V#13) 기준.')
legend(s)
chain(s, [
    ('HOST', '01', 'CSrvWork.ReadRequest() → Parsing()',
     'TCP 9911 수신, 전문 종류 판별'),
    ('HOST', '02', 'ParseOorR()',
     'O 전문 파싱. START_POS 가 9xx(크레인) 이면 출고'),
    ('HOST', '03', 'IsValidStartStation() / IsValidDestStation()',
     '출발 랙 셀과 도착 작업대 검증'),
    ('HOST', '04', 'frmMain.InsertJobMst()',
     'JOB_MST INSERT', 'JOB_STATUS = 99'),
    ('SCH', '05', 'AcceptNewJob()',
     '출고는 S/C 가 먼저 받는다', '99 → 20'),
    ('SCH', '06', 'DriveSC() → UpdateScData()',
     '해당 S/C 유휴(ONLINE·AUTO·OD_RQ_YN=\'N\') 확인 후 이송 명령', '20 → 25'),
    ('VEH', '07', 'VehThread.ConsumeCommands()',
     'S/C 에 출발 셀 / 도착 트랙 WRITE'),
    ('VEH', '08', 'VehThread.PollObservations()',
     'S/C 관측값을 SC_DATA_LGLS 에 반영'),
    ('SCH', '09', 'CompleteSC()',
     '출고 1차 완료 — 랙 셀 해제. 최종 아님', '25 → 29'),
    ('SCH', '10', 'LandScDrop()',
     '도착지 신호 OFF + HS_TRACK_NO 일치 → C/V 구간으로 인계', '29 → 15'),
    ('SCH', '11', 'DriveRGV() → UpdateRtvData()',
     'RGV 를 경유하는 경로일 때만. IsRtvBusyWithOwnJob() 이 1대 배타 관리', '15 → 35'),
    ('SCH', '12', 'CompleteRGVReal() → LandRgvDrop()',
     'RTV 완료 소비 후 착지 확인', '35 → 39 → 15'),
    ('SCH', '13', 'RequestArrivalTrackingWrite()',
     'C/V 도착 트랙에 작업번호 트래킹 기록. DriveCV() 는 10 전용이라 여기서는 돌지 않는다'),
    ('CV', '14', 'CvTrackingWrite()',
     'R 영역 작업번호를 PLC 로 WRITE — C/V 가 이 값으로 반송한다'),
    ('CV', '15', 'CvStatusScenario() → UpdateCvData()',
     'C/V 상태 READ → CV_DATA 미러'),
    ('CV', '16', 'CvEventCheck()',
     '출고대 도착 / 배출 완료 이벤트 ACK'),
    ('SCH', '17', 'CompleteCV()',
     'OD_RQ_YN=\'N\' + 도착 관측 후 배출 확인 → 출고 최종 상태', '15 → 19'),
    ('SCH', '18', 'ReportOutStationArrival()',
     '출고대 신호 ON → 상위 도착보고 대상(22)으로 표시'),
    ('HOST', '19', 'GetLoadArrivalReport()',
     '도착보고 전문 송신'),
    ('HOST', '20', 'GetJobCompleteReport(19) → IsJobExist(19)',
     'F 전문 조립'),
    ('HOST', '21', 'SendSock() → UpdateJobStatusTo(\'09\')',
     '완료보고 후 09. 실패 시 19 로 롤백', '19 → 09'),
])
footer(s, '출고 풀사이클')

# ══════════════════════════════════════════════════════════════════
# 5. IO_TASK 스케줄러 1사이클
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
head(s, 'IO_TASK 스케줄러 — Thread_Doing() 1사이클',
     '200ms 마다 아래 순서를 그대로 한 번 돈다. Drive* 는 대기 상태 + 유휴 설비에만 작용하는 멱등 함수다.')

groups = [
    ('접수', GREEN, 0.45, [
        ('AcceptNewJob()', '99 → 10(입고) / 20(출고)'),
    ]),
    ('① 구동 완료', BLUE, 0.45, [
        ('CompleteCV()', '15 → 19 (출고 최종) / 입고는 15 유지'),
        ('CompleteSC()', '25 → 29 (입고 최종 / 출고 1차)'),
        ('CompleteRGVReal()', '35 → 39 (RTV COMPLETE_RD 소비)'),
        ('CompleteRGVManual()', 'RTV 수동지시(9998) 완료 정리'),
        ('CompleteSCManual()', 'SC 수동지시(9999) 완료 정리'),
    ]),
    ('② 착지 처리', TEAL, 0.45, [
        ('LandRgvDrop()', '39 + HS_TRACK_NO 에 화물 → 15'),
        ('LandScDrop()', '출고 29 + 도착 신호 OFF → 15'),
    ]),
    ('③ 구동 지시', AMBER, 0.45, [
        ('SyncDualCvDirection()', '겸용대(C/V#11·#2) 방향 정합'),
        ('PromotePendingDirection()', '보류된 모드 전환(DIRW) 승격'),
        ('DriveCV()', '10 → 15'),
        ('DriveSC()', '입고 15 / 출고 20 → 25'),
        ('DriveRGV()', '15 → 35'),
    ]),
    ('④ 마무리 · 감시', RED, 0.45, [
        ('DeleteSemiFinished()', '반자동(11/12)은 19/29 에서 즉시 삭제 — 상위 보고 없음'),
        ('CheckStalledJobs()', '작업 체류(설비 무응답) 감시'),
        ('ReportOutStationArrival()', '출고대 신호 ON → 도착보고(22)'),
        ('MonitorAlarm() / CheckAlarm()', '설비 에러코드 로깅'),
        ('MarkErrorJobStatus()', '이중입고(54) / 공출고(58) → 작업상태 반영'),
        ('ResumeRedirectedJobs()', '재지정(07/06) → 새 셀로 재개 지시'),
        ('SyncDisplayTyp()', '표시용 작업구분 보강 (제어 미사용)'),
        ('Heartbeat(true)', '한 사이클 완주 → Client 상태표시줄 SCH 녹색'),
    ]),
]
PITCH = 0.212
y = 1.22
for title, col, x0, items in groups:
    h = 0.16 + PITCH * len(items)
    box(s, 0.45, y, 1.75, h, title, LIGHT, col, NAVY, 10.5)
    yy = y + 0.07
    for fn, desc in items:
        tb(s, 2.35, yy, 3.6, 0.22, fn, 9.5, True, NAVY)
        tb(s, 6.05, yy, 6.9, 0.22, desc, 9, False, GRAY)
        yy += PITCH
    y += h + 0.07

tb(s, 2.35, y + 0.06, 10.6,
   0.5, '※ 사이클이 예외로 깨지면 Heartbeat(false) 후 DB 재연결을 시도한다. Drive* 를 완료·착지 뒤에 한 번만 부르는 것은 '
        '방금 인계된 작업도 같은 호출이 잡기 때문이다 (2026-09-01 정리).', 9, False, GRAY)
footer(s, '스케줄러 1사이클')

# ══════════════════════════════════════════════════════════════════
# 6. WCS_TASK_CV 1사이클
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
head(s, 'WCS_TASK_CV — CvThread.Thread_Doing() 1사이클',
     '마스터 PLC 1소켓으로 15설비 슬롯을 순회한다. 한 바퀴가 약 16초라, 3초짜리 펄스는 DB 에서 놓칠 수 있다.')

box(s, 0.45, 1.30, 2.55, 0.62, ['ScanPendingWork()', ('쓰기 대기 1회 스캔', 8.5, False, GRAY)],
    PROC['CV'][1], AMBER, NAVY, 10.5)
arrow(s, 3.00, 1.61, 3.45, 1.61, LINE)
box(s, 3.45, 1.30, 2.55, 0.62, ['PreloadStatusBlocks()', ('상태 구간 일괄 조회 (왕복 45→3회)', 8.5, False, GRAY)],
    PROC['CV'][1], AMBER, NAVY, 10.5)
arrow(s, 6.00, 1.61, 6.45, 1.61, LINE)
box(s, 6.45, 1.30, 2.75, 0.62, ['지시 걸린 설비 우선 정렬', ('m_setPendCmd / Od / Trk', 8.5, False, GRAY)],
    LIGHT, GRAY, NAVY, 10.5)
arrow(s, 9.20, 1.61, 9.65, 1.61, LINE)
box(s, 9.65, 1.30, 3.25, 0.62, ['CvAlarmCheck(1)', ('알람 M비트는 전 C/V 공통 → 사이클당 1회', 8.5, False, GRAY)],
    PROC['CV'][1], AMBER, NAVY, 10.5)

tb(s, 0.45, 2.12, 12.4, 0.3, '슬롯(설비) 하나마다 아래 6단계를 순서대로 수행한다. 실패가 소켓 끊김이면 전체 재접속, 데이터 문제면 그 설비만 건너뛴다.',
   10, False, GRAY)

slot_steps = [
    ('GetFirstAddress(Idx)', '그 설비의 CV_DATA 로 읽을 첫 주소를 구한다. 실패하면 이 설비만 skip'),
    ('CvStatusScenario(Idx) / CvStatus(Idx)', 'PLC 상태 READ → UpdateCvData() 로 CV_DATA 미러'),
    ('CvChg_CMD_RQ_YN(Idx)', 'CMD 쓰기 지시를 PLC 로 WRITE (방향 전환 등)'),
    ('CvChg_OD_RQ_YN(Idx)', 'OD 쓰기 지시를 PLC 로 WRITE (이송 명령)'),
    ('CvEventCheck(Idx)', 'Load / Unload Complete 이벤트 ACK'),
    ('CvTrackingWrite(Idx)', 'R 영역 작업번호 트래킹 WRITE (EncodeJobNoR / GetRTrackingAddr)'),
]
y = 2.50
for i, (fn, desc) in enumerate(slot_steps, 1):
    box(s, 0.45, y, 0.42, 0.56, str(i), PROC['CV'][1], AMBER, NAVY, 11)
    box(s, 0.98, y, 4.35, 0.56, fn, WHITE, AMBER, NAVY, 10.5)
    tb(s, 5.50, y + 0.14, 7.4, 0.3, desc, 9.5, False, GRAY)
    if i < len(slot_steps):
        arrow(s, 0.66, y + 0.56, 0.66, y + 0.68, LINE, 1.25)
    y += 0.68

tb(s, 0.45, 6.62, 12.4, 0.4,
   '※ 연속 통신 실패가 임계를 넘으면 소켓을 닫고 스레드를 종료한다 — SYS_MAIN.Thread_Tick 이 재생성해 접속 루프를 다시 돈다 '
   '(2026-08-01, EQP_SIM 재기동 후 죽은 소켓에 34시간 매달린 사례 대응).', 9, False, GRAY)
footer(s, 'C/V 통신 1사이클')

# ══════════════════════════════════════════════════════════════════
# 7. WCS_TASK_HOST 송수신
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
head(s, 'WCS_TASK_HOST — 수신(9911) / 송신(9910) 함수',
     '수신은 CSrvWork(서버), 송신은 CCliWork(클라이언트)가 맡는다. 서로 다른 소켓이다.')

tb(s, 0.45, 1.22, 6.1, 0.3, '수신 — CSrvWork  (LocalPort 9911)', 13, True, NAVY)
recv = [
    ('ReadRequest()', 'CheckHeader() / CheckBody() 로 헤더·본문 길이 확인'),
    ('Parsing()', '전문 종류(Type) 판별 후 분기'),
    ('ParseOorR()', 'O = 작업지시 / R = 재지정. InsertJobMst() 로 작업 생성'),
    ('ParseModeChange()', 'M = 입출고 모드 전환 → SetCvDirection()'),
    ('ParseCancel()', '작업 취소 — 99 / 10 / 20 만 대상'),
    ('ParseArrived() / ParsePallet()', '도착 확인 / 팔레트 정보'),
    ('MakeResponse() → SendSock()', '응답 전문 회신 (정상 또는 반려 사유코드)'),
]
y = 1.58
for fn, desc in recv:
    box(s, 0.45, y, 2.95, 0.52, fn, PROC['HOST'][1], BLUE, NAVY, 10)
    tb(s, 3.52, y + 0.12, 3.0, 0.4, desc, 8.5, False, GRAY)
    if fn != recv[-1][0]:
        arrow(s, 1.92, y + 0.52, 1.92, y + 0.64, LINE, 1.25)
    y += 0.64

tb(s, 6.85, 1.22, 6.1, 0.3, '송신 — CCliWork  (RemotePort 9910)', 13, True, NAVY)
send = [
    ('GetSendData()', '아래를 매 주기 순서대로 호출한다', 0),
    ('GetStatusReport()', 'S = 상태 변경 시 + 30초 주기 상태보고', 0),
    ('GetErrorReport()', 'E = 설비 에러보고. IsEquip_ERROR_Modified() 로 변화 감지', 0),
    ('GetJobCompleteReport()', 'F = 완료보고. 대상 상태는 19(출고) / 29(입고)', 0),
    ('IsJobExist(nJobStatus)', '완료 대상 작업 조회', 1),
    ('UpdateJobStatusTo(lugg, \'09\')', '보고 성공 시 09. 실패하면 원래 상태로 롤백', 1),
    ('GetLoadArrivalReport()', '도착보고. GetDirOrder() / GetLuggOrder() 는 지시 하달용', 0),
]
y = 1.58
for i, (fn, desc, ind) in enumerate(send):
    dx = 0.38 * ind
    box(s, 6.85 + dx, y, 3.15 - dx, 0.52, fn, PROC['HOST'][1], BLUE, NAVY, 10)
    tb(s, 10.12, y + 0.12, 2.8, 0.4, desc, 8.5, False, GRAY)
    if ind:
        tb(s, 6.87, y + 0.13, 0.36, 0.26, '└', 11, True, LINE, PP_ALIGN.CENTER)
    elif i < len(send) - 1:
        arrow(s, 8.42, y + 0.52, 8.42, y + 0.64, LINE, 1.25)
    y += 0.64
tb(s, 6.85, y - 0.02, 6.05, 0.3, '└ 표시는 위 함수가 내부에서 부르는 함수다.', 8.5, False, LINE)

tb(s, 0.45, 6.20, 12.4, 0.7,
   ['※ 무게보고(U)와 빈파렛트 요청(P)은 이 현장이 재고를 관리하지 않아 발신을 막아 두었다 '
    '(GetWeightReport / GetEmptyPltRequest 는 호출되지 않는다).',
    '※ 시뮬 모드(chkSimMode)에서는 상태보고와 에러보고를 내지 않는다. 완료보고와 도착보고는 그대로 나간다.'],
   9, False, GRAY)
footer(s, '상위 통신 함수')

# ══════════════════════════════════════════════════════════════════
# 8. 예외 : 이중입고(54) / 공출고(58) 가 걸리는 지점
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
head(s, '예외가 걸리는 지점 — 이중입고(54) / 공출고(58)',
     '정상 흐름의 어느 함수에서 예외 경로로 빠지는지. 상세는 08_이중입고_공출고_시험_보고서 참조.')

box(s, 0.45, 1.30, 12.45, 0.55,
    ['정상 흐름 :  DriveSC() → VehThread.ConsumeCommands() → S/C 동작 → CompleteSC()  (25 → 29)'],
    LIGHT, GRAY, NAVY, 11)
arrow(s, 6.67, 1.85, 6.67, 2.15, RED, 1.75)
tb(s, 6.80, 1.86, 5.0, 0.28, '크레인이 셀에서 실패 → 에러코드 발생', 9.5, True, RED)

exc = [
    ('VehThread.PollObservations()', 'SC_DATA_LGLS.ERR_CODE_RD 에 0054 / 0058 기록. InsertEqpErrHis() 로 이력', AMBER),
    ('cThread_SCH.MarkErrorJobStatus()', '작업 상태를 에러 상태(08/07)로 표시. 재지정 상태(06/05)는 덮어쓰지 않는다', GREEN),
    ('CCliWork.GetErrorReport()', 'E 전문 보고. ErrorKind=1(이중입고) / 3(공출고). 작업번호는 실린 화물 기준', BLUE),
    ('CSrvWork.ParseOorR()  [R 분기]', '이중입고 : 같은 호기의 다른 셀로 R_Kind=1 재지정. 공출고 : 상위가 삭제한다', BLUE),
    ('cThread_SCH.ResumeRedirectedJobs()', '재지정(07/06) 작업을 새 셀로 재개 지시. 에러가 남아 있어도 지시를 기록한다', GREEN),
    ('설비 조작반 (사람)', '공출고는 재지정이 없어 새 지시가 안 들어간다 → 크레인 에러를 사람이 푼다', GRAY),
]
y = 2.15
for i, (fn, desc, col) in enumerate(exc):
    box(s, 0.45, y, 4.3, 0.58, fn, WHITE, col, NAVY, 10.5)
    tb(s, 4.92, y + 0.14, 8.0, 0.32, desc, 9.5, False, GRAY)
    if i < len(exc) - 1:
        arrow(s, 2.60, y + 0.58, 2.60, y + 0.70, LINE, 1.25)
    y += 0.70

tb(s, 0.45, 6.42, 12.45, 0.55,
   ['※ 공출고는 재지정이 아니라 삭제다 (IMS 담당자 확인). 상위가 작업을 삭제하고, 운전자가 [작업정보] 에서 같은 작업을 삭제한 뒤,',
    '   설비 조작반에서 크레인 에러를 해제한다. WCS 의 R_Kind=2(출고 재지정) 처리는 규격 지원으로만 남겨 두었다.'],
   9.5, False, RED)
footer(s, '예외 경로')

out = os.path.join(OUT, '09_함수단_흐름도.pptx')
prs.save(out)
print('saved', out, '|', page[0], 'slides')
