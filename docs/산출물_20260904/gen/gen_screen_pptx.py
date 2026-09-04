# -*- coding: utf-8 -*-
"""화면설계서 (pptx) - 2026-09-04 판 : 표지 / Document history / Contents / 화면별(캡처+설명)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from gen_common import OUT, SHOTS, DATE
from gen_screens import SCREENS

NAVY = RGBColor(0x1E, 0x27, 0x61); GRAY = RGBColor(0x59, 0x59, 0x59); LIGHT = RGBColor(0xEE, 0xF1, 0xF7); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]; FONT = '맑은 고딕'

def tb(slide, x, y, w, h, text, size=14, bold=False, color=GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = align
        if isinstance(ln, tuple): txt, b, sz, col = ln
        else: txt, b, sz, col = ln, bold, size, color
        r = p.add_run(); r.text = txt; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = col; r.font.name = FONT
    return s

page = [0]
def footer(slide):
    page[0] += 1
    tb(slide, 0.5, 7.0, 4, 0.35, 'LG 화학 1동 WCS  |  화면설계서', 10, False, GRAY)
    tb(slide, 12.2, 7.0, 0.8, 0.35, str(page[0]), 10, False, GRAY, PP_ALIGN.RIGHT)

def title(slide, text):
    tb(slide, 0.5, 0.35, 12.3, 0.7, text, 26, True, NAVY, anchor=MSO_ANCHOR.MIDDLE)

s = prs.slides.add_slide(blank); page[0] += 1
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tb(s, 1.0, 2.4, 11, 0.8, 'LG 화학 1동 자동창고', 22, False, RGBColor(0xCA, 0xDC, 0xFC))
tb(s, 1.0, 3.2, 11, 1.2, 'WCS 화면설계서', 48, True, WHITE)
tb(s, 1.0, 4.6, 11, 0.6, 'WCS Renewal (구 ECS 대체)   ·   Ver 1.1   ·   %s' % DATE, 16, False, RGBColor(0xCA, 0xDC, 0xFC))

def add_table(slide, x, y, w, rows, colw, size=11):
    n, m = len(rows), len(rows[0])
    shp = slide.shapes.add_table(n, m, Inches(x), Inches(y), Inches(w), Inches(0.4 * n)); t = shp.table
    for j, cw in enumerate(colw): t.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = t.cell(i, j); c.text = ''; p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = str(v); r.font.size = Pt(size); r.font.name = FONT
            c.fill.solid()
            if i == 0: c.fill.fore_color.rgb = NAVY; r.font.color.rgb = WHITE; r.font.bold = True
            else: c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT; r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    return t

s = prs.slides.add_slide(blank); title(s, 'Document history')
add_table(s, 0.7, 1.4, 11.9, [['Vers.', 'Date', 'Author', 'Approver', 'Notes'], ['V1.0', '2026-09-03', 'LGLS WCS Renewal', '', '최초 작성'], ['V1.1', DATE, 'LGLS WCS Renewal', '', '미사용 메뉴 정리, 캡처 갱신']], [1.2, 1.8, 2.6, 1.8, 4.5])
footer(s)

groups = []
for g, *_ in SCREENS:
    if g not in groups: groups.append(g)
s = prs.slides.add_slide(blank); title(s, 'Contents')
cols = 3; per = (len(groups) + cols - 1) // cols
for ci in range(cols):
    lines = []
    for g in groups[ci * per:(ci + 1) * per]:
        lines.append((g, True, 14, NAVY))
        for sc in [x for x in SCREENS if x[0] == g]: lines.append(('   ' + sc[1], False, 11, GRAY))
        lines.append(('', False, 6, GRAY))
    if lines: tb(s, 0.6 + ci * 4.15, 1.3, 4.0, 5.5, lines)
footer(s)

def screen_slide(g, name, sh, path, overview, notes):
    s = prs.slides.add_slide(blank); title(s, '%s - %s' % (g, name) if g != '메인' else name)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(8.4), Inches(5.5)); box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xFA); box.line.color.rgb = RGBColor(0xC8, 0xCE, 0xDA)
    p = os.path.join(SHOTS, sh) if sh else None
    if p and os.path.exists(p):
        im = Image.open(p); w, h = im.size; maxw, maxh = 8.2, 5.3
        sc = min(maxw / (w / 96.0), maxh / (h / 96.0)); dw, dh = w / 96.0 * sc, h / 96.0 * sc
        s.shapes.add_picture(p, Inches(0.6 + (8.4 - dw) / 2), Inches(1.3 + (5.5 - dh) / 2), Inches(dw), Inches(dh))
    else:
        pass
    lines = [('■ 메뉴 Path', True, 13, NAVY), (path, False, 11, GRAY), ('', False, 6, GRAY), ('■ 프로그램 개요', True, 13, NAVY), (overview, False, 11, GRAY)]
    if notes:
        lines += [('', False, 6, GRAY), ('■ 특기사항', True, 13, NAVY)] + [('· ' + n, False, 10, GRAY) for n in notes]
    tb(s, 9.3, 1.3, 3.6, 5.5, lines)
    footer(s)

for g, name, sh, path, overview, notes, ctrls in SCREENS:
    screen_slide(g, name, sh, path, overview, notes)
out = os.path.join(OUT, '06_WCS_화면설계서.pptx'); prs.save(out); print('saved', out, 'slides', len(prs.slides))
