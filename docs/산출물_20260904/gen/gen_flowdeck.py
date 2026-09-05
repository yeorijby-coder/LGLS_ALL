# -*- coding: utf-8 -*-
"""WCS 흐름도 (pptx) - 시스템 구성도 / 입고·출고 플로우차트 / 플로우 다이어그램"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUT = r'D:\project\LGLS\Renewal\docs\산출물_20260904'
DATE = '2026-09-05'
NAVY  = RGBColor(0x1E, 0x27, 0x61)
BLUE  = RGBColor(0x26, 0x6E, 0xBE)
ICE   = RGBColor(0xCA, 0xDC, 0xFC)
LIGHT = RGBColor(0xEE, 0xF1, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x59, 0x59, 0x59)
LINE  = RGBColor(0xB6, 0xBF, 0xD0)
GREEN = RGBColor(0x2E, 0x86, 0x4B)
AMBER = RGBColor(0xC8, 0x7A, 0x0E)
RED   = RGBColor(0xB0, 0x2E, 0x2E)
FONT  = '맑은 고딕'

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
page = [0]


def tb(slide, x, y, w, h, text, size=12, bold=False, color=GRAY,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = s.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = 0
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, tuple): txt, b, sz, col = ln
        else: txt, b, sz, col = ln, bold, size, color
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = col; r.font.name = FONT
    return s


def box(slide, x, y, w, h, text, fill=WHITE, edge=NAVY, txt=NAVY, size=11,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, bold=True):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text if isinstance(text, list) else [text]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        if isinstance(ln, tuple): t, sz, b, c = ln
        else: t, sz, b, c = ln, size, bold, txt
        r = p.add_run(); r.text = t
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = FONT
    return sh


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=1.5, dash=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width)
    el = c.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    tail = parse_xml('<a:tailEnd %s type="triangle" w="med" len="med"/>' % nsdecls('a'))
    el.append(tail)
    if dash:
        d = parse_xml('<a:prstDash %s val="dash"/>' % nsdecls('a'))
        el.insert(0, d)
    return c


def label(slide, x, y, w, text, size=9, color=GRAY, align=PP_ALIGN.CENTER):
    return tb(slide, x, y, w, 0.24, text, size, False, color, align)


def footer(slide, title):
    page[0] += 1
    tb(slide, 0.5, 7.02, 6, 0.3, 'LG 화학 1동 WCS  |  ' + title, 9, False, GRAY)
    tb(slide, 12.3, 7.02, 0.6, 0.3, str(page[0]), 9, False, GRAY, PP_ALIGN.RIGHT)


def head(slide, text, sub=None):
    tb(slide, 0.5, 0.3, 12.3, 0.5, text, 24, True, NAVY, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        tb(slide, 0.5, 0.82, 12.3, 0.3, sub, 11, False, GRAY)


# ══════════ 1. 표지 ══════════
s = prs.slides.add_slide(blank); page[0] += 1
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
tb(s, 1.0, 2.5, 11, 0.6, 'LG 화학 1동 자동창고', 20, False, ICE)
tb(s, 1.0, 3.15, 11, 1.1, 'WCS 시스템 구성도 및 업무 흐름도', 40, True, WHITE)
tb(s, 1.0, 4.5, 11, 0.5, 'WCS Renewal (구 ECS 대체)   ·   %s' % DATE, 14, False, ICE)

# ══════════ 2. 시스템 구성도 ══════════
s = prs.slides.add_slide(blank)
head(s, '시스템 구성도', '모든 프로그램은 DB(LGLS_MCS_IO)를 매개로 구동한다. 화살표는 데이터가 흐르는 방향이다.')

# 방사형 배치 : DB 를 가운데 두고 사방에 배치해 화살표가 박스를 지나가지 않게 한다.
box(s, 1.0, 2.20, 2.4, 0.80, ['상위 시스템', ('WMS / IMS', 10, False, GRAY)], LIGHT, NAVY, NAVY, 13)
box(s, 1.0, 3.60, 2.4, 0.85, ['상위 통신', ('WCS_TASK_HOST', 9, False, GRAY)], WHITE, BLUE, NAVY, 13)
box(s, 5.00, 1.60, 2.4, 0.85, ['스케줄러', ('IO_TASK', 9, False, GRAY)], WHITE, BLUE, NAVY, 13)
box(s, 4.85, 3.35, 2.7, 1.35, ['DB', ('Microsoft SQL Server', 9, False, WHITE), ('LGLS_MCS_IO', 11, True, WHITE)],
    NAVY, NAVY, WHITE, 16, MSO_SHAPE.FLOWCHART_MAGNETIC_DISK)
box(s, 5.00, 5.50, 2.4, 0.85, ['설비 통신', ('WCS_TASK_CV', 9, False, GRAY)], WHITE, BLUE, NAVY, 13)
box(s, 9.20, 3.60, 2.6, 0.85, ['운전 화면', ('Ecs.exe (Client)', 9, False, GRAY)], WHITE, BLUE, NAVY, 13)
box(s, 9.20, 5.50, 2.6, 0.85, ['설비 / PLC', ('XGT · C/V · S/C · RGV', 9, False, GRAY)], LIGHT, NAVY, NAVY, 13)

arrow(s, 2.20, 3.00, 2.20, 3.60)            # 상위 시스템 → 상위 통신
arrow(s, 2.20, 3.60, 2.20, 3.00)            # 상위 통신 → 상위 시스템
arrow(s, 3.40, 4.02, 4.85, 4.02)            # 상위 통신 → DB
arrow(s, 6.20, 2.45, 6.20, 3.35)            # 스케줄러 → DB
arrow(s, 6.35, 3.35, 6.35, 2.45)            # DB → 스케줄러
arrow(s, 6.20, 5.50, 6.20, 4.70)            # 설비 통신 → DB
arrow(s, 6.35, 4.70, 6.35, 5.50)            # DB → 설비 통신
arrow(s, 9.20, 4.02, 7.55, 4.02)            # 운전 화면 → DB
arrow(s, 7.40, 5.92, 9.20, 5.92, GREEN)     # 설비 통신 → 설비
arrow(s, 9.20, 6.12, 7.40, 6.12, GREEN)     # 설비 → 설비 통신

label(s, 1.0, 4.50, 2.4, 'O/R/M/S/E/F 전문 · TCP 8001·8002')
label(s, 7.55, 6.15, 1.65, 'XGT 1소켓 · 2004')

tb(s, 8.30, 1.30, 4.5, 2.0, [
    ('■ 각 프로그램의 역할', True, 12, NAVY),
    ('· 상위 통신 : WMS/IMS 전문을 작업(JOB_MST)으로 바꾸고', False, 10.5, GRAY),
    ('  상태·완료·에러를 되보고한다.', False, 10.5, GRAY),
    ('· 스케줄러 : 작업 상태를 진행시키고 설비 지시(_OD)를 쓴다.', False, 10.5, GRAY),
    ('· 설비 통신 : PLC 를 읽어 _RD 에 쓰고, 지시를 PLC 에 전송한다.', False, 10.5, GRAY),
    ('· 운전 화면 : DB 를 보여주고 수동 지시·설정을 DB 에 기록한다.', False, 10.5, GRAY),
    ('  설비와 직접 통신하지 않는다.', False, 10.5, GRAY),
])
tb(s, 0.6, 5.55, 3.9, 1.2, [
    ('■ 한 줄 요약', True, 12, NAVY),
    ('프로그램끼리 직접 부르지 않는다.', False, 10.5, GRAY),
    ('모두 DB 를 통해 주고받으므로,', False, 10.5, GRAY),
    ('하나가 죽어도 지시와 상태는 DB 에 남는다.', False, 10.5, GRAY),
])
footer(s, '시스템 구성도')

# ══════════ 3. 입고 플로우차트 ══════════
s = prs.slides.add_slide(blank)
head(s, '입고 플로우차트', '작업 상태(JOB_STATUS)를 따라가는 판단 흐름. 로직1(124) 기준이며 로직2(130)·로직3(122)은 출발 트랙만 다르다.')

X = 5.6; W = 2.15
ys = [1.25, 1.95, 2.65, 3.35, 4.05, 4.75, 5.45, 6.15]
steps = [
    ('시작 : 입고대에 파렛트 적재', LIGHT, '', MSO_SHAPE.ROUNDED_RECTANGLE),
    ('상위 O 전문 수신\n작업 생성 (99)', WHITE, '상위 통신', MSO_SHAPE.RECTANGLE),
    ('입고대 CV 구동지시 (10 → 15)', WHITE, '스케줄러', MSO_SHAPE.RECTANGLE),
    ('RGV 반송 (15 → 35)\n픽업 트랙 → 크레인 통로', WHITE, '스케줄러 · RGV', MSO_SHAPE.RECTANGLE),
    ('RGV 하역 완료 (39)', WHITE, '설비 통신', MSO_SHAPE.RECTANGLE),
    ('도착 트랙에 작업번호 기록 (39 → 15)', WHITE, '스케줄러', MSO_SHAPE.RECTANGLE),
    ('크레인 반송 (15 → 25 → 29)\nS/C측 트랙 → 랙 셀', WHITE, '스케줄러 · 크레인', MSO_SHAPE.RECTANGLE),
    ('완료 (09) · 상위 F 전문 보고', LIGHT, '상위 통신', MSO_SHAPE.ROUNDED_RECTANGLE),
]
for i, (t, f, who, shp) in enumerate(steps):
    h = 0.55 if '\n' not in t else 0.58
    box(s, X, ys[i], W + 1.4, h, t.split('\n'), f, NAVY, NAVY, 11, shp)
    if who:
        tb(s, X - 2.0, ys[i] + 0.13, 1.85, 0.3, who, 10, True, BLUE, PP_ALIGN.RIGHT)
    if i < len(steps) - 1:
        arrow(s, X + (W + 1.4) / 2, ys[i] + h, X + (W + 1.4) / 2, ys[i + 1])

# 우측 : 보류 조건
tb(s, 9.55, 1.25, 3.3, 0.35, '■ RGV 지시 전 확인 조건', 12, True, NAVY)
for i, t in enumerate(['① 출발 트랙 : RGV 출발 H/S 신호 ON',
                       '② 도착 트랙 : RGV 도착 H/S 신호 ON',
                       '③ RGV : 대기(IDLE) 상태',
                       '④ 도착 트랙에 착지 기록 대기(39) 작업 없음']):
    tb(s, 9.55, 1.68 + i * 0.42, 3.3, 0.4, t, 10, False, GRAY)
box(s, 9.55, 3.55, 3.3, 0.75,
    ['하나라도 아니면 지시하지 않고', '알람 이력에 사유를 남긴다'], LIGHT, AMBER, AMBER, 10)

tb(s, 9.55, 4.6, 3.3, 0.35, '■ 상태 코드', 12, True, NAVY)
for i, t in enumerate(['99 신규   10 CV 구동대기   15 CV 구동중',
                       '30 RGV 대기   35 RGV 구동중   39 RGV 완료',
                       '20 크레인 대기   25 구동중   29 완료',
                       '09 완료(이력 이관)']):
    tb(s, 9.55, 5.0 + i * 0.33, 3.3, 0.3, t, 10, False, GRAY)
footer(s, '입고 플로우차트')

# ══════════ 4. 입고 플로우 다이어그램 ══════════
LANES = [('WMS / IMS', 1.30), ('상위 통신', 2.15), ('스케줄러', 3.00),
         ('설비 통신', 3.85), ('설비 / PLC', 4.70)]

def swimlanes(slide, note):
    for name, y in LANES:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.72))
        bar.fill.solid(); bar.fill.fore_color.rgb = LIGHT
        bar.line.color.rgb = LINE; bar.line.width = Pt(0.75); bar.shadow.inherit = False
        bar.text_frame.text = ''
        tb(slide, 0.6, y + 0.22, 1.5, 0.3, name, 11, True, NAVY)
    tb(slide, 0.5, 5.6, 12.3, 1.3, note, 11, False, GRAY)

def step(slide, x, lane, w, text, fill=WHITE, edge=BLUE, size=9.5):
    y = LANES[lane][1] + 0.09
    return box(slide, x, y, w, 0.54, text if isinstance(text, list) else [text],
               fill, edge, NAVY, size, MSO_SHAPE.ROUNDED_RECTANGLE, True)

def hop(slide, x, l1, l2, color=BLUE):
    y1 = LANES[l1][1] + (0.72 if l2 > l1 else 0.09)
    y2 = LANES[l2][1] + (0.09 if l2 > l1 else 0.72)
    arrow(slide, x, y1, x, y2, color)

s = prs.slides.add_slide(blank)
head(s, '입고 플로우 다이어그램', '누가 무엇을 하는지 시간 순서(왼쪽 → 오른쪽)로 본 흐름')
swimlanes(s, ['· 스케줄러는 설비에 직접 말하지 않는다. 지시를 DB(_OD)에 쓰면 설비 통신이 PLC 로 전송한다.',
              '· 설비가 올린 완료 신호는 설비 통신이 관측해 Ack 를 돌려주고 DB(_RD)에 남긴다. 스케줄러는 그 값을 보고 상태를 올린다.'])
xs = [1.95, 3.55, 5.15, 6.75, 8.35, 9.95, 11.35]
step(s, xs[0], 0, 1.45, ['O 전문', '(작업 지시)'], LIGHT)
hop(s, xs[0] + 0.72, 0, 1)
step(s, xs[0], 1, 1.45, ['작업 생성', 'JOB_MST 99'])
hop(s, xs[0] + 0.72, 1, 2)
step(s, xs[1], 2, 1.45, ['CV 지시', '10 → 15'])
hop(s, xs[1] + 0.72, 2, 3)
step(s, xs[1], 3, 1.45, ['PLC 쓰기', '트래킹/목적지'])
hop(s, xs[1] + 0.72, 3, 4)
step(s, xs[1], 4, 1.45, ['컨베이어 이송', '→ RGV 픽업'], LIGHT, NAVY)
step(s, xs[2], 2, 1.45, ['RGV 지시', '15 → 35'])
hop(s, xs[2] + 0.72, 2, 3)
step(s, xs[2], 3, 1.45, ['반송지시 전송'])
hop(s, xs[2] + 0.72, 3, 4)
step(s, xs[2], 4, 1.45, ['RGV 반송', '통로로 하역'], LIGHT, NAVY)
step(s, xs[3], 4, 1.45, ['하역 완료 신호'], LIGHT, NAVY)
hop(s, xs[3] + 0.72, 4, 3, GREEN)
step(s, xs[3], 3, 1.45, ['완료 관측 · Ack', 'COMPLETE_RD'], WHITE, GREEN)
hop(s, xs[3] + 0.72, 3, 2, GREEN)
step(s, xs[3], 2, 1.45, ['35 → 39 → 15'], WHITE, GREEN)
step(s, xs[4], 2, 1.45, ['크레인 지시', '15 → 25'])
hop(s, xs[4] + 0.72, 2, 3)
step(s, xs[4], 3, 1.45, ['반송지시 전송'])
hop(s, xs[4] + 0.72, 3, 4)
step(s, xs[4], 4, 1.45, ['크레인 적재', '셀에 저장'], LIGHT, NAVY)
step(s, xs[5], 4, 1.25, ['완료 신호'], LIGHT, NAVY)
hop(s, xs[5] + 0.62, 4, 3, GREEN)
step(s, xs[5], 3, 1.25, ['완료 관측 · Ack'], WHITE, GREEN)
hop(s, xs[5] + 0.62, 3, 2, GREEN)
step(s, xs[5], 2, 1.25, ['25 → 29 → 09'], WHITE, GREEN)
hop(s, xs[5] + 0.62, 2, 1, GREEN)
step(s, xs[5], 1, 1.25, ['F 전문 보고'], WHITE, GREEN)
hop(s, xs[5] + 0.62, 1, 0, GREEN)
step(s, xs[6], 0, 1.35, ['완료 확인'], LIGHT)
arrow(s, xs[5] + 1.25, LANES[0][1] + 0.36, xs[6], LANES[0][1] + 0.36, GREEN)
footer(s, '입고 플로우 다이어그램')

# ══════════ 5. 출고 플로우차트 ══════════
s = prs.slides.add_slide(blank)
head(s, '출고 플로우차트', '입고의 역순이다. 크레인이 먼저 움직이고 마지막에 출고대에서 지게차가 가져간다.')
steps = [
    ('시작 : 상위 O 전문(출발 셀 · 도착 작업대)', LIGHT, '상위 통신', MSO_SHAPE.ROUNDED_RECTANGLE),
    ('작업 생성 (99) → 크레인 구동대기 (20)', WHITE, '스케줄러', MSO_SHAPE.RECTANGLE),
    ('크레인 반송 (20 → 25 → 29)\n랙 셀 → S/C측 트랙', WHITE, '스케줄러 · 크레인', MSO_SHAPE.RECTANGLE),
    ('통로 CV 가 RGV측 트랙으로 이송 (15)', WHITE, '설비', MSO_SHAPE.RECTANGLE),
    ('RGV 반송 (35 → 39)\n통로 → 출고대 픽업 트랙', WHITE, '스케줄러 · RGV', MSO_SHAPE.RECTANGLE),
    ('출고대 진입 · 출고대 신호 ON (15 → 19)', WHITE, '설비 통신', MSO_SHAPE.RECTANGLE),
    ('지게차 반출 → 화물감지 OFF', WHITE, '작업자', MSO_SHAPE.RECTANGLE),
    ('완료 (09) · 상위 F 전문 보고', LIGHT, '상위 통신', MSO_SHAPE.ROUNDED_RECTANGLE),
]
for i, (t, f, who, shp) in enumerate(steps):
    h = 0.55 if '\n' not in t else 0.58
    box(s, X, ys[i], W + 1.4, h, t.split('\n'), f, NAVY, NAVY, 11, shp)
    if who:
        tb(s, X - 2.0, ys[i] + 0.13, 1.85, 0.3, who, 10, True, BLUE, PP_ALIGN.RIGHT)
    if i < len(steps) - 1:
        arrow(s, X + (W + 1.4) / 2, ys[i] + h, X + (W + 1.4) / 2, ys[i + 1])

tb(s, 9.55, 1.25, 3.3, 0.35, '■ 출고 작업대', 12, True, NAVY)
for i, t in enumerate(['126 원부자재 불출대 (C/V#13, 로직1)',
                       '129 피킹 작업대 (C/V#14, 로직2)',
                       '122 입출고 겸용대 (C/V#11, 로직3)']):
    tb(s, 9.55, 1.68 + i * 0.34, 3.3, 0.32, t, 10, False, GRAY)
box(s, 9.55, 2.9, 3.3, 0.95,
    ['겸용대(122)는 출고 모드일 때만', '출고 화물이 나올 수 있다.', '모드 전환은 상위 M 전문(WMS → WCS)'],
    LIGHT, AMBER, AMBER, 10)
tb(s, 9.55, 4.15, 3.3, 0.35, '■ 상태 코드', 12, True, NAVY)
for i, t in enumerate(['99 신규   20 크레인 대기   25 구동중   29 완료',
                       '15 CV 구동중   35 RGV 구동중   39 RGV 완료',
                       '19 CV 구동완료(출고대 도착)',
                       '09 완료(이력 이관)']):
    tb(s, 9.55, 4.55 + i * 0.33, 3.3, 0.3, t, 10, False, GRAY)
footer(s, '출고 플로우차트')

# ══════════ 6. 출고 플로우 다이어그램 ══════════
s = prs.slides.add_slide(blank)
head(s, '출고 플로우 다이어그램', '누가 무엇을 하는지 시간 순서(왼쪽 → 오른쪽)로 본 흐름')
swimlanes(s, ['· 출고는 크레인이 먼저 움직인다. 크레인이 통로에 내려놓으면 컨베이어와 RGV 가 출고대까지 이어 나른다.',
              '· 마지막 완료는 사람이 만든다. 지게차가 화물을 가져가 출고대 화물감지가 꺼지면 그때 완료(19 → 09)로 처리한다.'])
step(s, xs[0], 0, 1.45, ['O 전문', '(출고 지시)'], LIGHT)
hop(s, xs[0] + 0.72, 0, 1)
step(s, xs[0], 1, 1.45, ['작업 생성', 'JOB_MST 99'])
hop(s, xs[0] + 0.72, 1, 2)
step(s, xs[1], 2, 1.45, ['크레인 지시', '20 → 25'])
hop(s, xs[1] + 0.72, 2, 3)
step(s, xs[1], 3, 1.45, ['반송지시 전송'])
hop(s, xs[1] + 0.72, 3, 4)
step(s, xs[1], 4, 1.45, ['크레인 출고', '셀 → 트랙'], LIGHT, NAVY)
step(s, xs[2], 4, 1.45, ['완료 신호'], LIGHT, NAVY)
hop(s, xs[2] + 0.72, 4, 3, GREEN)
step(s, xs[2], 3, 1.45, ['완료 관측 · Ack'], WHITE, GREEN)
hop(s, xs[2] + 0.72, 3, 2, GREEN)
step(s, xs[2], 2, 1.45, ['25 → 29 → 15'], WHITE, GREEN)
step(s, xs[3], 2, 1.45, ['RGV 지시', '15 → 35'])
hop(s, xs[3] + 0.72, 2, 3)
step(s, xs[3], 3, 1.45, ['반송지시 전송'])
hop(s, xs[3] + 0.72, 3, 4)
step(s, xs[3], 4, 1.45, ['RGV 반송', '→ 출고대'], LIGHT, NAVY)
step(s, xs[4], 4, 1.45, ['출고대 신호 ON'], LIGHT, NAVY)
hop(s, xs[4] + 0.72, 4, 3, GREEN)
step(s, xs[4], 3, 1.45, ['신호 관측'], WHITE, GREEN)
hop(s, xs[4] + 0.72, 3, 2, GREEN)
step(s, xs[4], 2, 1.45, ['39 → 15 → 19'], WHITE, GREEN)
step(s, xs[5], 4, 1.25, ['지게차 반출', '감지 OFF'], LIGHT, NAVY)
hop(s, xs[5] + 0.62, 4, 3, GREEN)
step(s, xs[5], 3, 1.25, ['해제 관측'], WHITE, GREEN)
hop(s, xs[5] + 0.62, 3, 2, GREEN)
step(s, xs[5], 2, 1.25, ['19 → 09'], WHITE, GREEN)
hop(s, xs[5] + 0.62, 2, 1, GREEN)
step(s, xs[5], 1, 1.25, ['F 전문 보고'], WHITE, GREEN)
hop(s, xs[5] + 0.62, 1, 0, GREEN)
step(s, xs[6], 0, 1.35, ['완료 확인'], LIGHT)
arrow(s, xs[5] + 1.25, LANES[0][1] + 0.36, xs[6], LANES[0][1] + 0.36, GREEN)
footer(s, '출고 플로우 다이어그램')

# ══════════ 7. 피킹 작업대 ══════════
s = prs.slides.add_slide(blank)
head(s, '피킹 작업대 흐름 (102)', '129(출고)와 130(재입고)은 물리적으로 한 자리다. IMS 는 둘을 하나의 작업대 102 로 본다.')
box(s, 1.0, 1.5, 2.5, 0.9, ['출고 흐름', '랙 → 129 도착'], WHITE, NAVY, NAVY, 12)
box(s, 4.2, 1.5, 2.5, 0.9, ['작업자 피킹', '필요 수량만 꺼냄'], LIGHT, AMBER, AMBER, 12)
box(s, 7.4, 1.5, 2.5, 0.9, ['PLC 인계', '#30 ON → #29 OFF'], WHITE, GREEN, GREEN, 12)
box(s, 10.6, 1.5, 2.2, 0.9, ['재입고 흐름', '130 → 랙'], WHITE, NAVY, NAVY, 12)
arrow(s, 3.5, 1.95, 4.2, 1.95); arrow(s, 6.7, 1.95, 7.4, 1.95); arrow(s, 9.9, 1.95, 10.6, 1.95)

tb(s, 1.0, 2.85, 11.8, 0.4, '■ PLC 사양(시나리오 슬라이드 13)에 적힌 인계 순서', 13, True, NAVY)
seq = ['① 화물이 #29 에 도착 → W.O 비트 ON, 트래킹 유지',
       '② PLC Unload Request #2 ON → WCS Ack ON   (작업자 피킹 구간)',
       '③ 반출 완료 → PLC Unload Complete #2 ON → WCS Ack ON → 양쪽 OFF',
       '④ Pallet Exist #30 ON   (아직 #29 도 ON)',
       '⑤ Pallet Exist #29 OFF   ※ 사양 주석 : "출고 영역 #29 bit 만 지움"']
for i, t in enumerate(seq):
    tb(s, 1.2, 3.3 + i * 0.38, 11.5, 0.36, t, 11, False, GRAY)

box(s, 1.0, 5.35, 11.8, 1.35,
    [('확인 필요 (PLC 담당자)', 12, True, AMBER),
     ('· 작업자 완료 버튼은 없다고 확인됨 → Unload Complete #2 는 PLC 타이머가 올리는 것으로 보인다.', 11, False, GRAY),
     ('· #30 재하 ON 조건에 WCS Ack(M0933)가 들어가는지, 타이머 단독인지 확인이 필요하다.', 11, False, GRAY),
     ('· #30 트래킹(작업번호)은 WCS 가 기록한다 — 재입고는 상위의 새 지시로 시작한다.', 11, False, GRAY)],
    LIGHT, AMBER, GRAY, 11, MSO_SHAPE.ROUNDED_RECTANGLE, False)
footer(s, '피킹 작업대 흐름')

# ══════════ 8. 예외 처리 ══════════
s = prs.slides.add_slide(blank)
head(s, '예외 처리 흐름', '이중입고는 상위가 위치를 재지정하고, 공출고는 재지정 없이 양쪽에서 작업을 삭제한다.')

tb(s, 0.7, 1.20, 5.9, 0.4, '■ 이중입고 (에러 54) — 입고 목적 셀에 이미 화물', 13, True, NAVY)
tb(s, 0.9, 1.62, 5.7, 0.32, '→ 상위가 R 전문으로 다른 셀을 내려 준다 (재지정 있음)', 10.5, True, BLUE)
seq1 = ['① 크레인이 셀에 넣지 못하고 에러 54 발생 (화물은 실은 채)',
        '② WCS 가 E 전문 보고 (ErrorKind=1, 작업번호 · 도착 셀)',
        '③ 상위가 같은 호기의 다른 셀로 R 전문(R_Kind=1) 재지정',
        '④ WCS 가 작업을 재지정 상태로 바꾸고 새 셀로 크레인 재지시',
        '⑤ 크레인이 새 지시를 받아 정상 입고 완료']
for i2, t in enumerate(seq1):
    tb(s, 0.9, 2.05 + i2 * 0.40, 5.7, 0.38, t, 11, False, GRAY)

tb(s, 6.95, 1.20, 5.9, 0.4, '■ 공출고 (에러 58) — 출고 출발 셀이 비어 있음', 13, True, NAVY)
tb(s, 7.15, 1.62, 5.7, 0.32, '→ 재지정이 없다. 양쪽에서 작업을 삭제한다', 10.5, True, AMBER)
seq2 = ['① 크레인이 셀에서 집지 못하고 에러 58 발생 (빈 포크)',
        '② WCS 가 E 전문 보고 (ErrorKind=3, 작업번호 · 출발 셀)',
        '③ 상위(IMS)가 그 작업을 삭제한다 (담당자 확인 후)',
        '④ 운전자가 WCS [작업정보] 에서 같은 작업을 삭제 (확인 후)',
        '⑤ 운전자가 설비 조작반에서 크레인 에러를 해제한다',
        '⑥ 재고가 실제로 있는 셀로 상위가 새 출고를 지시']
for i2, t in enumerate(seq2):
    tb(s, 7.15, 2.05 + i2 * 0.40, 5.7, 0.38, t, 11, False, GRAY)

box(s, 0.7, 4.28, 12.15, 0.95,
    [('설비 에러 해제는 운전 화면의 기능이 아니다', 11.5, True, NAVY),
     ('구 ECS 에도 S/C · RGV 에러를 해제하는 화면 기능은 없다. 설비에서 알람이 해제되면 설비가 AlarmResetReport 로 알리고 ECS 는 Ack 만 한다.', 10.5, False, GRAY),
     ('이중입고는 재지정 지시가 들어가면서 설비가 스스로 풀지만, 공출고는 재지정이 없어 조작반에서 사람이 풀어야 한다.', 10.5, False, GRAY),
     ('(참고) 구 ECS 소스에서 R 전문은 ACK 만 하고 실제 위치 변경 코드가 없다. 재지정 동작은 현행 WCS 에서 구현했다.', 10.5, False, GRAY)],
    LIGHT, NAVY, GRAY, 10.5, MSO_SHAPE.ROUNDED_RECTANGLE, False)

tb(s, 0.7, 5.32, 12.1, 0.4, '■ 그 밖의 정체와 조치', 13, True, NAVY)
rows = [['상황', '조치'],
        ['같은 상태에 오래 머묾 (정체 알림창)', '설비 대화상자에서 상태 확인 → 지시 재전송 / 지시 삭제'],
        ['RGV 가 지시를 받지 않음', 'RTV 상태창 [지시 삭제] 로 요청 플래그까지 비운다'],
        ['설비 완료 신호가 오지 않아 상태가 남음', '[환경설정] > [시간기반 자동처리] 를 선택하면 경과시간으로 완료를 추정한다'],
        ['상태바 EQUIP / HOST / SCH 빨강', '해당 서버 프로그램 재기동 (지시는 DB 에 남아 이어서 진행된다)']]
shp = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(5.62), Inches(12.1), Inches(0.24 * len(rows)))
t = shp.table
t.columns[0].width = Inches(4.3); t.columns[1].width = Inches(7.8)
for i2, row in enumerate(rows):
    for j2, v in enumerate(row):
        c = t.cell(i2, j2); c.text = ''
        p2 = c.text_frame.paragraphs[0]; r = p2.add_run(); r.text = v
        r.font.size = Pt(9); r.font.name = FONT
        c.fill.solid()
        if i2 == 0:
            c.fill.fore_color.rgb = NAVY; r.font.color.rgb = WHITE; r.font.bold = True
        else:
            c.fill.fore_color.rgb = WHITE if i2 % 2 else LIGHT
            r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
footer(s, '예외 처리 흐름')

out = os.path.join(OUT, '07_WCS_구성도_및_흐름도.pptx')
prs.save(out)
print('saved', out, 'slides', len(prs.slides))
